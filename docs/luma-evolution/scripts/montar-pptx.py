#!/usr/bin/env python3
"""docs/luma-evolution/scripts/montar-pptx.py — empacota os slides em .pptx.

    python3 docs/luma-evolution/scripts/montar-pptx.py

Cada slide entra como uma imagem 1920x1080 ocupando a página inteira, e as notas do
apresentador vão pro campo de notas do PowerPoint.

POR QUE IMAGEM E NÃO CAIXAS DE TEXTO: o layout já foi desenhado e validado em HTML.
Reimplementá-lo em formas do PowerPoint criaria uma segunda versão do mesmo visual pra
divergir da primeira — e é justamente aí que nascem texto cortado e imagem deformada.
Empacotando o PNG validado, o .pptx mostra exatamente o que foi revisado.

Custo assumido: o texto do .pptx não é selecionável nem editável no PowerPoint. Quem
precisar mexer no conteúdo mexe no HTML e roda `node scripts/tudo.js` de novo — que é o
fluxo que mantém as três saídas iguais.
"""
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    sys.exit("falta o python-pptx: pip install python-pptx")

BASE = Path(__file__).resolve().parent.parent
PNGS = BASE / "presentation" / "slides-png"
NOTAS = BASE / "presentation" / "luma-evolution-speaker-notes.md"
SAIDA = BASE / "presentation" / "luma-evolution.pptx"

# 16:9 em EMU (1 polegada = 914400 EMU). 13,333 x 7,5 polegadas = 1920x1080 a 144 dpi.
LARG, ALT = Emu(12192000), Emu(6858000)


def notas_por_slide():
    """Lê as notas do apresentador e devolve {numero_do_slide: texto}."""
    if not NOTAS.exists():
        return {}
    blocos = re.split(r"^## Slide (\d+)\s*$", NOTAS.read_text(encoding="utf-8"), flags=re.M)
    return {int(blocos[i]): blocos[i + 1].strip() for i in range(1, len(blocos) - 1, 2)}


def main():
    imagens = sorted(PNGS.glob("slide-*.png"))
    if not imagens:
        sys.exit(f"nenhum PNG em {PNGS} — rode antes: node docs/luma-evolution/scripts/publicar.js")

    notas = notas_por_slide()
    prs = Presentation()
    prs.slide_width, prs.slide_height = LARG, ALT
    vazio = prs.slide_layouts[6]  # layout em branco: sem placeholder pra atrapalhar

    for img in imagens:
        n = int(re.search(r"(\d+)", img.stem).group(1))
        s = prs.slides.add_slide(vazio)
        s.shapes.add_picture(str(img), 0, 0, width=LARG, height=ALT)
        if n in notas:
            s.notes_slide.notes_text_frame.text = notas[n]

    prs.save(SAIDA)
    com_notas = sum(1 for i in imagens if int(re.search(r"(\d+)", i.stem).group(1)) in notas)
    print(f"PPTX: {SAIDA.name} — {len(imagens)} slides, {com_notas} com notas do apresentador")


if __name__ == "__main__":
    main()
